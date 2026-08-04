"""Format-aware document extraction with source-location metadata."""
from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Optional



@dataclass(frozen=True)
class ExtractedUnit:
    text: str
    page_number: Optional[int] = None
    slide_number: Optional[int] = None
    section_heading: Optional[str] = None


@dataclass(frozen=True)
class ExtractedDocument:
    file_type: str
    units: list[ExtractedUnit]
    page_count: Optional[int] = None
    slide_count: Optional[int] = None

    @property
    def text(self) -> str:
        return "\n\n".join(unit.text for unit in self.units if unit.text.strip())


def _pdf(path: Path) -> ExtractedDocument:
    from PyPDF2 import PdfReader

    reader = PdfReader(str(path))
    units = [
        ExtractedUnit(text=(page.extract_text() or "").strip(), page_number=index)
        for index, page in enumerate(reader.pages, start=1)
    ]
    return ExtractedDocument("pdf", units, page_count=len(reader.pages))


def _pptx(path: Path) -> ExtractedDocument:
    from pptx import Presentation

    presentation = Presentation(str(path))
    units: list[ExtractedUnit] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        title = slide.shapes.title.text.strip() if slide.shapes.title and slide.shapes.title.text.strip() else None
        notes = ""
        try:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        except Exception:
            # Notes are optional and python-pptx support varies by version.
            pass
        if notes:
            texts.append(notes)
        units.append(ExtractedUnit("\n".join(texts), slide_number=index, section_heading=title))
    return ExtractedDocument("pptx", units, slide_count=len(presentation.slides))


def _docx(path: Path) -> ExtractedDocument:
    import docx

    document = docx.Document(str(path))
    units: list[ExtractedUnit] = []
    heading: Optional[str] = None
    body: list[str] = []

    def flush() -> None:
        if body:
            units.append(ExtractedUnit("\n".join(body), section_heading=heading))
            body.clear()

    for paragraph in document.paragraphs:
        value = paragraph.text.strip()
        if not value:
            continue
        if paragraph.style and paragraph.style.name.startswith("Heading"):
            flush()
            heading = value
        else:
            body.append(value)
    flush()
    return ExtractedDocument("docx", units)


def _txt(path: Path) -> ExtractedDocument:
    return ExtractedDocument("txt", [ExtractedUnit(path.read_text(encoding="utf-8"))])


def extract_document(file_path: str) -> ExtractedDocument:
    """Extract a supported upload, preserving its natural structural units."""
    path = Path(file_path)
    extension = path.suffix.lower()
    if extension == ".pdf":
        return _pdf(path)
    if extension in {".ppt", ".pptx"}:
        return _pptx(path)
    if extension in {".doc", ".docx"}:
        return _docx(path)
    if extension == ".txt":
        return _txt(path)
    raise ValueError(f"Unsupported file type: {extension}")
